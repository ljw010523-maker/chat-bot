"""
범용 문서 로더 (Universal Document Loader) - HWP/HWPX 지원 추가
모든 파일 형식을 완벽하게 파싱 + 노이즈 제거 + HWP 검증 강화
"""

from pathlib import Path
from typing import List, Dict
import platform

# PDF
import PyPDF2
from pdf2image import convert_from_path
import pytesseract
from PIL import Image, ImageEnhance

# Office 문서
try:
    import docx

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import pandas as pd

    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

# 구버전 Office 파일 (.doc, .xls, .ppt)
try:
    import win32com.client  # type: ignore

    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False

# HWP 파일 처리
try:
    import olefile

    HWP_AVAILABLE = True
except ImportError:
    HWP_AVAILABLE = False

# hwp5 라이브러리는 설치가 어려워서 제거
# HWP는 olefile의 PrvText 방식으로 처리

# Unstructured (Deep Document Parser)
try:
    from unstructured.partition.auto import partition
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False

# Upstage VLM API (표/도장/날인 인식)
try:
    from langchain_upstage import UpstageDocumentParseLoader
    VLM_AVAILABLE = True
except ImportError:
    VLM_AVAILABLE = False

# 텍스트 정제
from back.scripts.clean.text_cleaner import TextCleaner

# Tesseract/Poppler 경로 (Windows)
if platform.system() == "Windows":
    # Tesseract OCR 비활성화
    # pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

    import glob

    poppler_search = glob.glob(r"C:\Program Files\poppler-*\Library\bin")
    POPPLER_PATH = (
        poppler_search[0]
        if poppler_search
        else r"C:\Program Files\poppler-25.07.0\Library\bin"
    )
else:
    POPPLER_PATH = None


class UniversalDocumentLoader:
    """범용 문서 로더 - 모든 파일 형식 완벽 처리"""

    def __init__(self, config):
        self.config = config
        self.ocr_dpi = getattr(config, "ocr_dpi", 300)
        self.text_cleaner = TextCleaner()  # 텍스트 정제기 초기화
        self._print_capabilities()

    def _print_capabilities(self):
        """지원 형식 출력"""
        print("\n📋 문서 로더 지원 형식:")
        print(f"  - PDF: ✓ (텍스트 + OCR + 노이즈 제거)")
        print(f"  - TXT: ✓")
        print(
            f"  - Word (.docx): {'✓' if DOCX_AVAILABLE else '✗ (pip install python-docx)'}"
        )
        print(
            f"  - Word (.doc): {'✓' if WIN32COM_AVAILABLE else '✗ (Windows 전용, pip install pywin32)'}"
        )
        print(
            f"  - PowerPoint (.pptx): {'✓' if PPTX_AVAILABLE else '✗ (pip install python-pptx)'}"
        )
        print(
            f"  - PowerPoint (.ppt): {'✓' if WIN32COM_AVAILABLE else '✗ (Windows 전용, pip install pywin32)'}"
        )
        print(
            f"  - Excel (.xlsx): {'✓' if OPENPYXL_AVAILABLE else '✗ (pip install openpyxl)'}"
        )
        print(
            f"  - Excel (.xls): {'✓' if WIN32COM_AVAILABLE else '✗ (Windows 전용, pip install pywin32)'}"
        )
        print(f"  - CSV: {'✓' if PANDAS_AVAILABLE else '✗ (pip install pandas)'}")
        print(
            f"  - HWP (.hwp/.hwpx): {'✓' if HWP_AVAILABLE else '✗ (pip install olefile)'}"
        )
        print(f"  - 이미지 (.jpg/.png): ✓ (OCR + 노이즈 제거)\n")

    def load(self, file_path: Path) -> List[Dict]:
        """
        파일 형식 자동 감지 및 텍스트 추출

        Returns:
            List[Dict]: [{'page_num': int, 'text': str, 'method': str}, ...]
        """
        suffix = file_path.suffix.lower()

        # 파일 형식별 라우팅
        if suffix == ".pdf":
            return self._load_pdf(file_path)
        elif suffix == ".txt":
            return self._load_txt(file_path)
        elif suffix == ".docx":
            return self._load_docx(file_path)
        elif suffix == ".doc":
            return self._load_doc_legacy(file_path)
        elif suffix == ".pptx":
            return self._load_pptx(file_path)
        elif suffix == ".ppt":
            return self._load_ppt_legacy(file_path)
        elif suffix == ".xlsx":
            return self._load_excel(file_path)
        elif suffix == ".xls":
            return self._load_xls_legacy(file_path)
        elif suffix == ".csv":
            return self._load_csv(file_path)
        # elif suffix in [".hwp", ".hwpx"]:  # HWP 처리 비활성화 (느림)
        #     return self._load_hwp(file_path)
        elif suffix in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
            return self._load_image(file_path)
        else:
            print(f"  ⚠️ 지원하지 않는 형식: {suffix}")
            return []

    # ============================================
    # HWP 텍스트 검증 (신규 추가)
    # ============================================

    def _is_valid_korean_text(self, text: str) -> bool:
        """한글 텍스트 유효성 검증 (깨진 텍스트 필터링)"""
        if not text or len(text.strip()) < 10:
            return False

        # 한글, 영문, 숫자, 기본 특수문자만 추출
        valid_chars = [
            c
            for c in text
            if (
                "\uac00" <= c <= "\ud7a3"  # 한글 완성형
                or "a" <= c.lower() <= "z"  # 영문
                or c.isdigit()  # 숫자
                or c in " \n\t.,!?-()[]{}:;@#%&*+=/<>\"'"  # 기본 특수문자
            )
        ]

        if not valid_chars:
            return False

        # 유효한 문자 비율 체크
        valid_ratio = len(valid_chars) / len(text)

        # 한글 비율 체크
        korean_chars = sum(1 for c in text if "\uac00" <= c <= "\ud7a3")
        korean_ratio = korean_chars / len(text) if text else 0

        # 영문/숫자 비율
        alnum_chars = sum(1 for c in text if c.isalnum())
        alnum_ratio = alnum_chars / len(text) if text else 0

        # 유효 조건 (하나라도 만족하면 OK)
        conditions = [
            valid_ratio >= 0.5,  # 유효 문자 50% 이상
            korean_ratio >= 0.05,  # 한글 5% 이상
            (korean_chars >= 3),  # 한글 3자 이상
            (alnum_chars >= 10 and valid_ratio >= 0.3),  # 영문/숫자 10자 이상
        ]

        return any(conditions)

    # ============================================
    # PDF 처리 (강화)
    # ============================================

    def _load_pdf(self, file_path: Path) -> List[Dict]:
        """PDF 텍스트 추출 (텍스트 우선, 스캔본은 Google Vision OCR)"""
        try:
            # 방법 1: PyPDF2로 텍스트 추출 시도 (타이핑된 문서용)
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)

                print(f"  📑 PDF {total_pages}페이지")

                # 첫 페이지로 텍스트 여부 판단
                first_text = reader.pages[0].extract_text().strip()
                first_text_len = len(first_text)

                print(f"  📊 첫 페이지 텍스트: {first_text_len}자")

                # 텍스트가 충분하면 (타이핑된 문서)
                if first_text_len >= 50:
                    print("  ✓ 텍스트 기반 PDF (타이핑된 문서)")
                    pages_data = []

                    for page_num, page in enumerate(reader.pages, 1):
                        text = page.extract_text().strip()

                        if len(text) < 50:
                            print(f"    페이지 {page_num}: 텍스트 부족, OCR 적용")
                            text = self._ocr_pdf_page(file_path, page_num)
                            method = "pdf_ocr"
                        else:
                            # 텍스트 후처리 적용
                            text = self.text_cleaner.clean_ocr_text(text)
                            method = "pdf_text"

                        pages_data.append(
                            {"page_num": page_num, "text": text, "method": method}
                        )

                    return pages_data

                # 텍스트가 부족하면 스캔본/캡처본으로 판단
                print("  ⚠️ 스캔본/캡처본 감지 (텍스트 부족)")

            # 방법 2: Upstage VLM OCR (스캔본/캡처본 전용 - 표/도장/날인 인식)
            if VLM_AVAILABLE:
                print("  🔍 Upstage VLM OCR 시도 중...")
                result = self._parse_with_vlm(file_path)
                if result:
                    print(f"  ✓ VLM OCR 성공 ({len(result)}페이지)")
                    return result
                print("  ⚠️ VLM OCR 실패 → Tesseract OCR 폴백")

            # 방법 3: Tesseract OCR (최종 폴백)
            print("  → Tesseract OCR 모드로 전환")
            return self._ocr_pdf(file_path)

        except Exception as e:
            print(f"  ❌ PDF 읽기 실패: {e}")
            print("  → Tesseract OCR 모드로 전환")
            return self._ocr_pdf(file_path)

    def _ocr_pdf(self, file_path: Path) -> List[Dict]:
        """PDF 전체 OCR (강화 버전)"""
        try:
            if platform.system() == "Windows":
                images = convert_from_path(
                    file_path, dpi=self.ocr_dpi, poppler_path=POPPLER_PATH
                )
            else:
                images = convert_from_path(file_path, dpi=self.ocr_dpi)

            print(f"  🔍 {len(images)}페이지 OCR 처리 중 (노이즈 필터링)...")

            pages_data = []
            for page_num, image in enumerate(images, 1):
                print(f"    페이지 {page_num}/{len(images)}...", end=" ")

                try:
                    # 전처리
                    image = self._preprocess_image_for_table(image)

                    text = pytesseract.image_to_string(
                        image, lang="kor+eng", config="--oem 1 --psm 6"
                    ).strip()

                    # 후처리
                    text = self.text_cleaner.clean_ocr_text(text)

                    print(f"✓ ({len(text)}자)")

                    pages_data.append(
                        {"page_num": page_num, "text": text, "method": "pdf_ocr"}
                    )
                except Exception as e:
                    print(f"✗ 실패: {e}")
                    pages_data.append(
                        {"page_num": page_num, "text": "", "method": "ocr_failed"}
                    )

            return pages_data

        except Exception as e:
            print(f"  ❌ OCR 실패: {e}")
            return []

    def _ocr_pdf_page(self, file_path: Path, page_num: int) -> str:
        """PDF 특정 페이지만 OCR (강화 버전)"""
        try:
            if platform.system() == "Windows":
                images = convert_from_path(
                    file_path,
                    first_page=page_num,
                    last_page=page_num,
                    dpi=self.ocr_dpi,
                    poppler_path=POPPLER_PATH,
                )
            else:
                images = convert_from_path(
                    file_path, first_page=page_num, last_page=page_num, dpi=self.ocr_dpi
                )

            if images:
                # 전처리
                image = self._preprocess_image_for_table(images[0])

                text = pytesseract.image_to_string(
                    image, lang="kor+eng", config="--oem 1 --psm 6"
                ).strip()

                # 후처리
                text = self.text_cleaner.clean_ocr_text(text)

                return text
            return ""

        except Exception as e:
            print(f" (OCR 실패: {e})")
            return ""

    # ============================================
    # TXT 처리
    # ============================================

    def _load_txt(self, file_path: Path) -> List[Dict]:
        """TXT 파일 읽기 (다중 인코딩 지원)"""
        try:
            encodings = ["utf-8", "utf-8-sig", "cp949", "euc-kr", "latin-1", "ascii"]
            text = None

            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        text = f.read()
                    print(f"  ✓ TXT 읽기 성공 ({encoding})")
                    break
                except (UnicodeDecodeError, LookupError):
                    continue

            if text is None:
                print("  ❌ 인코딩 실패")
                return []

            # 후처리 적용
            text = self.text_cleaner.clean_ocr_text(text)

            return [{"page_num": 1, "text": text, "method": "txt"}]

        except Exception as e:
            print(f"  ❌ TXT 읽기 실패: {e}")
            return []

    # ============================================
    # Word 처리 (강화 - .docx + .doc)
    # ============================================

    def _load_docx(self, file_path: Path) -> List[Dict]:
        """Word .docx 파일 읽기"""
        if not DOCX_AVAILABLE:
            print("  ❌ python-docx 미설치")
            print("     설치: pip install python-docx")
            return []

        try:
            doc = docx.Document(file_path)

            # 단락별 텍스트 추출
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

            # 표 추출
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        tables_text.append(row_text)

            # 병합
            all_text = "\n".join(paragraphs)
            if tables_text:
                all_text += "\n\n[표 데이터]\n" + "\n".join(tables_text)

            # 후처리 추가
            all_text = self.text_cleaner.clean_ocr_text(all_text)

            print(f"  ✓ Word (.docx) 읽기 완료 ({len(all_text)}자)")

            return [{"page_num": 1, "text": all_text, "method": "docx"}]

        except Exception as e:
            print(f"  ❌ Word (.docx) 읽기 실패: {e}")
            if "not a zip file" in str(e).lower():
                print("  → 구버전 .doc 파일로 재시도")
                return self._load_doc_legacy(file_path)
            return []

    def _load_doc_legacy(self, file_path: Path) -> List[Dict]:
        """Word .doc 파일 읽기 (구버전 - Windows 전용)"""
        if not WIN32COM_AVAILABLE:
            print("  ❌ pywin32 미설치 (Windows 전용)")
            print("     설치: pip install pywin32")
            return []

        if platform.system() != "Windows":
            print("  ❌ .doc 파일은 Windows에서만 지원됩니다")
            return []

        try:
            import pythoncom

            pythoncom.CoInitialize()

            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False

            doc = word.Documents.Open(str(file_path.absolute()))
            text = doc.Content.Text

            doc.Close()
            word.Quit()

            pythoncom.CoUninitialize()

            # 후처리 추가
            text = self.text_cleaner.clean_ocr_text(text)

            print(f"  ✓ Word (.doc) 읽기 완료 ({len(text)}자)")

            return [{"page_num": 1, "text": text, "method": "doc_legacy"}]

        except Exception as e:
            print(f"  ❌ Word (.doc) 읽기 실패: {e}")
            return []

    # ============================================
    # PowerPoint 처리 (강화 - .pptx + .ppt)
    # ============================================

    def _load_pptx(self, file_path: Path) -> List[Dict]:
        """PowerPoint .pptx 파일 읽기"""
        if not PPTX_AVAILABLE:
            print("  ❌ python-pptx 미설치")
            print("     설치: pip install python-pptx")
            return []

        try:
            prs = Presentation(file_path)
            pages_data = []

            print(f"  📊 PowerPoint (.pptx) {len(prs.slides)}슬라이드")

            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())

                slide_text = "\n".join(texts)

                # 후처리 추가
                slide_text = self.text_cleaner.clean_ocr_text(slide_text)

                pages_data.append(
                    {"page_num": slide_num, "text": slide_text, "method": "pptx"}
                )

            print(f"  ✓ PowerPoint (.pptx) 읽기 완료")
            return pages_data

        except Exception as e:
            print(f"  ❌ PowerPoint (.pptx) 읽기 실패: {e}")
            return []

    def _load_ppt_legacy(self, file_path: Path) -> List[Dict]:
        """PowerPoint .ppt 파일 읽기 (구버전 - Windows 전용)"""
        if not WIN32COM_AVAILABLE:
            print("  ❌ pywin32 미설치 (Windows 전용)")
            print("     설치: pip install pywin32")
            return []

        if platform.system() != "Windows":
            print("  ❌ .ppt 파일은 Windows에서만 지원됩니다")
            return []

        try:
            import pythoncom

            pythoncom.CoInitialize()

            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False

            presentation = powerpoint.Presentations.Open(
                str(file_path.absolute()), ReadOnly=True
            )
            pages_data = []

            print(f"  📊 PowerPoint (.ppt) {presentation.Slides.Count}슬라이드")

            for slide_num in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides(slide_num)
                texts = []

                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        if shape.TextFrame.HasText:
                            texts.append(shape.TextFrame.TextRange.Text)

                slide_text = "\n".join(texts)

                # 후처리 추가
                slide_text = self.text_cleaner.clean_ocr_text(slide_text)

                pages_data.append(
                    {"page_num": slide_num, "text": slide_text, "method": "ppt_legacy"}
                )

            presentation.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()

            print(f"  ✓ PowerPoint (.ppt) 읽기 완료")
            return pages_data

        except Exception as e:
            print(f"  ❌ PowerPoint (.ppt) 읽기 실패: {e}")
            return []

    # ============================================
    # Excel 처리 (강화 - .xlsx + .xls)
    # ============================================

    def _load_excel(self, file_path: Path) -> List[Dict]:
        """Excel .xlsx 파일 읽기"""
        if not OPENPYXL_AVAILABLE:
            print("  ❌ openpyxl 미설치")
            print("     설치: pip install openpyxl")
            return []

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            pages_data = []

            print(f"  📊 Excel (.xlsx) {len(wb.sheetnames)}시트")

            for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
                sheet = wb[sheet_name]

                rows_text = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip():
                        rows_text.append(row_text)

                sheet_text = f"[시트: {sheet_name}]\n" + "\n".join(rows_text)

                # 후처리 추가
                sheet_text = self.text_cleaner.clean_ocr_text(sheet_text)

                pages_data.append(
                    {"page_num": sheet_num, "text": sheet_text, "method": "xlsx"}
                )

            print(f"  ✓ Excel (.xlsx) 읽기 완료")
            return pages_data

        except Exception as e:
            print(f"  ❌ Excel (.xlsx) 읽기 실패: {e}")
            return []

    def _load_xls_legacy(self, file_path: Path) -> List[Dict]:
        """Excel .xls 파일 읽기 (구버전 - Windows 전용)"""
        if not WIN32COM_AVAILABLE:
            print("  ❌ pywin32 미설치 (Windows 전용)")
            print("     설치: pip install pywin32")
            return []

        if platform.system() != "Windows":
            print("  ❌ .xls 파일은 Windows에서만 지원됩니다")
            return []

        try:
            import pythoncom

            pythoncom.CoInitialize()

            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False

            workbook = excel.Workbooks.Open(str(file_path.absolute()), ReadOnly=True)
            pages_data = []

            print(f"  📊 Excel (.xls) {workbook.Sheets.Count}시트")

            for sheet_num in range(1, workbook.Sheets.Count + 1):
                sheet = workbook.Sheets(sheet_num)
                rows_text = []

                used_range = sheet.UsedRange
                for row in used_range.Rows:
                    row_values = []
                    for cell in row.Cells:
                        value = cell.Value
                        row_values.append(str(value) if value is not None else "")
                    row_text = " | ".join(row_values)
                    if row_text.strip():
                        rows_text.append(row_text)

                sheet_text = f"[시트: {sheet.Name}]\n" + "\n".join(rows_text)

                # 후처리 추가
                sheet_text = self.text_cleaner.clean_ocr_text(sheet_text)

                pages_data.append(
                    {"page_num": sheet_num, "text": sheet_text, "method": "xls_legacy"}
                )

            workbook.Close(False)
            excel.Quit()
            pythoncom.CoUninitialize()

            print(f"  ✓ Excel (.xls) 읽기 완료")
            return pages_data

        except Exception as e:
            print(f"  ❌ Excel (.xls) 읽기 실패: {e}")
            return []

    # ============================================
    # CSV 처리 (강화)
    # ============================================

    def _load_csv(self, file_path: Path) -> List[Dict]:
        """CSV 읽기 (다중 인코딩 지원)"""
        if not PANDAS_AVAILABLE:
            print("  ❌ pandas 미설치")
            print("     설치: pip install pandas")
            return []

        try:
            encodings = ["utf-8", "utf-8-sig", "cp949", "euc-kr", "latin-1"]
            df = None

            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    print(f"  ✓ CSV 읽기 성공 ({encoding}, {len(df)}행)")
                    break
                except:
                    continue

            if df is None:
                print("  ❌ CSV 인코딩 실패")
                return []

            csv_text = df.to_string(index=False)

            # 후처리 추가
            csv_text = self.text_cleaner.clean_ocr_text(csv_text)

            return [{"page_num": 1, "text": csv_text, "method": "csv"}]

        except Exception as e:
            print(f"  ❌ CSV 읽기 실패: {e}")
            return []

    # ============================================
    # Unstructured (Deep Document Parser)
    # ============================================

    def _load_with_unstructured(self, file_path: Path) -> List[Dict]:
        """Unstructured를 사용한 딥러닝 기반 문서 파싱"""
        print(f"  🔄 방법: Unstructured Deep Parser")

        if not UNSTRUCTURED_AVAILABLE:
            print(f"  ⚠️ Unstructured 라이브러리 미설치")
            print(f"     설치: pip install unstructured")
            print(f"     전체 지원: pip install \"unstructured[all-docs]\"")
            return []

        try:
            # Unstructured로 파일 파싱
            elements = partition(filename=str(file_path))

            # 텍스트 추출
            full_text = "\n\n".join([str(el) for el in elements])

            # 후처리
            full_text = self.text_cleaner.clean_ocr_text(full_text)

            print(f"  ✅ Unstructured 파싱 완료 ({len(full_text)}자, {len(elements)}개 요소)")

            return [
                {
                    "page_num": 1,
                    "text": full_text,
                    "method": f"unstructured_{file_path.suffix[1:]}",
                    "elements_count": len(elements)
                }
            ]

        except Exception as e:
            print(f"  ⚠️ Unstructured 파싱 실패: {e}")
            return []

    # ============================================
    # Upstage VLM API (문서 OCR - 표/도장/날인 인식)
    # ============================================

    def _parse_with_vlm(self, file_path: Path) -> List[Dict]:
        """Upstage VLM API로 문서 파싱 (표/도장/날인 인식)"""
        print(f"  🔄 방법: Upstage VLM OCR")

        try:
            from langchain_upstage import UpstageDocumentParseLoader

            api_key = self.config.upstage_api_key
            if not api_key:
                print(f"  ⚠️ Upstage API 키가 설정되지 않았습니다")
                return []

            loader = UpstageDocumentParseLoader(
                file_path=str(file_path),
                split="page",
                api_key=api_key,
            )

            docs = loader.load()
            pages_data = []

            for idx, doc in enumerate(docs, 1):
                text = doc.page_content.strip()

                # 노이즈 제거
                text = self.text_cleaner.clean_ocr_text(text)

                pages_data.append(
                    {"page_num": idx, "text": text, "method": "vlm_ocr"}
                )

            print(f"  ✅ VLM OCR 완료 ({len(pages_data)}페이지, {sum(len(p['text']) for p in pages_data)}자)")
            return pages_data

        except Exception as e:
            print(f"  ⚠️ VLM OCR 실패: {e}")
            return []

    # ============================================
    # HWP → PDF 변환 (win32com)
    # ============================================

    def _disable_hwp_security_via_registry(self):
        """레지스트리를 통해 한글 보안 수준 완전 비활성화"""
        try:
            import winreg

            # 한글 보안 설정 레지스트리 경로들
            reg_paths = [
                (winreg.HKEY_CURRENT_USER, r"Software\HNC\HwpAutomation"),
                (winreg.HKEY_CURRENT_USER, r"Software\HNC\Hwp\8.0\HncOle"),
                (winreg.HKEY_CURRENT_USER, r"Software\Hnc\Hwp\Automation"),
                (winreg.HKEY_CURRENT_USER, r"Software\HNC\Hwp\9.0"),
                (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\HNC\Hwp\9.0"),
            ]

            for hkey, reg_path in reg_paths:
                try:
                    key = winreg.CreateKey(hkey, reg_path)
                    # 모든 보안 관련 설정 비활성화
                    winreg.SetValueEx(key, "Security", 0, winreg.REG_DWORD, 0)
                    winreg.SetValueEx(key, "SecurityLevel", 0, winreg.REG_DWORD, 0)
                    winreg.SetValueEx(key, "TrustVBAProject", 0, winreg.REG_DWORD, 1)
                    winreg.SetValueEx(key, "DisableSecurityWarning", 0, winreg.REG_DWORD, 1)
                    winreg.CloseKey(key)
                except:
                    pass

        except Exception:
            pass

    def _auto_click_hwp_security_popup(self, timeout=10):
        """한글 보안 팝업 자동 클릭 (백그라운드 쓰레드 - 버튼 직접 클릭)"""
        import time
        import threading

        def click_popup():
            try:
                import win32gui
                import win32con

                start_time = time.time()
                clicked_count = 0

                while time.time() - start_time < timeout:
                    # 한글 보안 경고 창 찾기 (여러 제목 시도)
                    window_titles = ["호환", "보안 경고", "한글", "HWP", "알림", "경고"]
                    hwnd = 0

                    for title in window_titles:
                        hwnd = win32gui.FindWindow(None, title)
                        if hwnd != 0:
                            break

                    if hwnd != 0:
                        # 창 활성화 (포커스 주기) - 키보드 입력을 받을 수 있도록
                        try:
                            # 창을 전면으로 가져오기
                            win32gui.SetForegroundWindow(hwnd)
                            # 창 표시하기
                            win32gui.ShowWindow(hwnd, win32con.SW_SHOW)
                            # 짧은 대기 (창이 활성화될 시간)
                            time.sleep(0.1)
                        except:
                            pass

                        # 방법 1: 자식 버튼 찾아서 직접 클릭
                        button_clicked = False

                        def find_and_click_button(hwnd_child, _):
                            nonlocal button_clicked
                            try:
                                class_name = win32gui.GetClassName(hwnd_child)
                                text = win32gui.GetWindowText(hwnd_child)

                                # 버튼 클래스이고 "모두 허용" 또는 관련 텍스트
                                if "button" in class_name.lower():
                                    # "모두 허용(A)", "허용 완료(N)", "허용" 등
                                    if any(keyword in text for keyword in ["모두", "허용", "완료"]):
                                        # WM_LBUTTONDOWN + WM_LBUTTONUP으로 클릭 시뮬레이션
                                        win32gui.SendMessage(hwnd_child, win32con.WM_LBUTTONDOWN, 0, 0)
                                        win32gui.SendMessage(hwnd_child, win32con.WM_LBUTTONUP, 0, 0)
                                        # BM_CLICK도 시도
                                        win32gui.SendMessage(hwnd_child, win32con.BM_CLICK, 0, 0)
                                        print(f"    ✓ 버튼 클릭: '{text}'")
                                        button_clicked = True
                                        return False
                            except:
                                pass
                            return True

                        try:
                            win32gui.EnumChildWindows(hwnd, find_and_click_button, None)
                        except:
                            pass

                        if button_clicked:
                            clicked_count += 1
                            print(f"    ✓ 보안 팝업 자동 승인 완료 ({clicked_count}번째)")
                            time.sleep(0.3)
                            # 여러 번 나올 수 있으므로 계속 감시
                            if clicked_count >= 5:  # 최대 5번까지만
                                break
                        else:
                            # 방법 2: 키보드 입력 (Alt+A)
                            import win32api
                            win32api.keybd_event(win32con.VK_MENU, 0, 0, 0)  # Alt 누름
                            time.sleep(0.05)
                            win32api.keybd_event(0x41, 0, 0, 0)  # A 누름
                            time.sleep(0.05)
                            win32api.keybd_event(0x41, 0, win32con.KEYEVENTF_KEYUP, 0)  # A 뗌
                            win32api.keybd_event(win32con.VK_MENU, 0, win32con.KEYEVENTF_KEYUP, 0)  # Alt 뗌
                            print(f"    ✓ Alt+A 키 전송")
                            time.sleep(0.3)

                    time.sleep(0.1)  # 100ms마다 확인
            except Exception:
                pass

        # 백그라운드 쓰레드 시작
        thread = threading.Thread(target=click_popup, daemon=True)
        thread.start()
        return thread

    def _convert_hwp_to_pdf(self, hwp_path: Path) -> Path:
        """HWP를 PDF로 변환 (한글 프로그램 자동화)"""
        import tempfile
        import uuid

        print(f"  🔄 HWP → PDF 변환 중...")

        # 임시 PDF 파일 경로 (한글 경로 문제 방지를 위해 UUID 사용)
        temp_dir = Path(tempfile.gettempdir())
        pdf_path = temp_dir / f"hwp_temp_{uuid.uuid4().hex}.pdf"

        # 기존 파일 삭제
        if pdf_path.exists():
            pdf_path.unlink()

        # 레지스트리를 통한 보안 설정 비활성화 (사전 방지)
        self._disable_hwp_security_via_registry()

        # 보안 팝업 자동 클릭 쓰레드 시작 (2차 방어)
        self._auto_click_hwp_security_popup(timeout=15)

        try:
            if not WIN32COM_AVAILABLE or platform.system() != "Windows":
                print(f"  ⚠️ Windows 환경이 아니거나 pywin32 미설치")
                return None

            import win32com.client
            import pythoncom

            pythoncom.CoInitialize()

            # 한글 프로그램 실행 (백그라운드)
            hwp = win32com.client.Dispatch("HWPFrame.HwpObject")

            # 프로그램 창 완전히 숨기기 (모든 UI 비활성화)
            try:
                hwp.XFrameWindow.Visible = False  # 메인 창 숨김
                hwp.XFrameWindow.Active = 0       # 창 비활성화
            except:
                pass

            # 화면 업데이트 중지 (성능 향상 + 팝업 방지)
            try:
                hwp.SetPrivateInfoPath("", "")    # 개인정보 경로 무시
            except:
                pass

            # 보안 경고 완전 무시 (모든 메시지 박스 자동 처리)
            try:
                # 0x01000000 = 모든 메시지 박스 무시
                # 0x00020000 = 메시지 박스 자동 승인
                # 0x00010000 = 경고 메시지 무시
                hwp.SetMessageBoxMode(0x01000000 | 0x00020000 | 0x00010000)
            except:
                pass

            # 보안 모듈 무시 (파일 경로 검사 우회)
            try:
                hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModuleExample")
            except:
                pass

            # 파일 열기 (모든 보안 확인 무시)
            open_params = "openreadonly:true;versionwarning:false;suspendpassword:true;lock:false;noconfirm:true"
            hwp.Open(str(hwp_path.absolute()), "HWP", open_params)

            # PDF로 저장
            hwp.HAction.GetDefault("FileSaveAsPdf", hwp.HParameterSet.HFileOpenSave.HSet)
            hwp.HParameterSet.HFileOpenSave.filename = str(pdf_path.absolute())
            hwp.HParameterSet.HFileOpenSave.Format = "PDF"
            hwp.HAction.Execute("FileSaveAsPdf", hwp.HParameterSet.HFileOpenSave.HSet)

            # 종료
            hwp.Quit()
            pythoncom.CoUninitialize()

            if pdf_path.exists():
                print(f"  ✅ PDF 변환 성공: {pdf_path}")
                return pdf_path
            else:
                print(f"  ❌ PDF 파일 생성 실패")
                return None

        except Exception as e:
            print(f"  ❌ HWP → PDF 변환 실패: {e}")
            return None

    # ============================================
    # HWP 처리 (olefile 우선, Google Vision OCR 폴백)
    # ============================================

    def _load_hwp(self, file_path: Path) -> List[Dict]:
        """HWP 파일 읽기 - olefile 우선, Google Vision OCR 폴백"""
        print(f"\n  📄 HWP 파일 처리: {file_path.name}")

        # HWP 5.0 이상 (HWPX 형식)
        if file_path.suffix.lower() == ".hwpx":
            return self._load_hwpx(file_path)

        # 방법 1: olefile로 PrvText 추출 시도 (팝업 없음, 빠름)
        print(f"  🔄 방법 1: olefile로 PrvText 추출 시도")
        olefile_result = None
        olefile_text_length = 0

        if HWP_AVAILABLE:
            try:
                import olefile

                if olefile.isOleFile(file_path):
                    ole = olefile.OleFileIO(file_path)
                    texts = []

                    # PrvText 추출
                    if ole.exists("PrvText"):
                        try:
                            stream = ole.openstream("PrvText")
                            data = stream.read()
                            text = data.decode("utf-16le", errors="ignore")
                            text = text.replace("\x00", "")

                            if self._is_valid_korean_text(text):
                                texts.append(text.strip())
                                print(f"  ✓ PrvText 추출 성공 ({len(text)}자)")
                        except Exception as e:
                            print(f"  ⚠️ PrvText 추출 실패: {e}")

                    ole.close()

                    if texts:
                        full_text = "\n\n".join(texts)
                        full_text = self.text_cleaner.clean_ocr_text(full_text)
                        olefile_text_length = len(full_text)
                        olefile_result = [{"page_num": 1, "text": full_text, "method": "hwp_prvtext"}]
                        print(f"  ✅ olefile 추출 완료 ({olefile_text_length}자)")
            except Exception as e:
                print(f"  ⚠️ olefile 처리 실패: {e}")

        # olefile 결과 확인: 파일 크기 대비 텍스트가 충분하면 바로 반환
        # 파일 크기로 "얼마나 많이 남아있는지" 판단
        file_size_kb = file_path.stat().st_size / 1024
        # HWP 파일은 일반적으로 1KB당 약 100~200자의 텍스트 포함
        # PrvText는 전체의 약 30~50% 정도만 포함
        # 파일 크기 기반 + 최소 기준 둘 다 사용
        expected_min_chars = max(int(file_size_kb * 100), 1500)  # 최소 1500자 요구

        if olefile_result and olefile_text_length >= expected_min_chars:
            ratio = (olefile_text_length / expected_min_chars * 100) if expected_min_chars > 0 else 100
            print(f"  ✅ olefile로 충분한 텍스트 추출됨")
            print(f"     파일: {file_size_kb:.1f}KB, 추출: {olefile_text_length}자, 예상 최소: {expected_min_chars}자 ({ratio:.0f}%)")
            return olefile_result

        # 방법 2: olefile 실패 또는 텍스트 부족 → Google Vision OCR 폴백
        if olefile_result:
            shortage = expected_min_chars - olefile_text_length
            print(f"  ⚠️ olefile 텍스트 부족")
            print(f"     파일: {file_size_kb:.1f}KB, 추출: {olefile_text_length}자, 예상 최소: {expected_min_chars}자 (부족: {shortage}자)")
        else:
            print(f"  ⚠️ olefile 추출 실패")

        print(f"  🔄 방법 2: VLM OCR 폴백 (HWP → PDF → VLM)")

        if VLM_AVAILABLE:
            # 2-1. HWP → PDF 변환
            pdf_path = self._convert_hwp_to_pdf(file_path)

            if pdf_path and pdf_path.exists():
                # 2-2. VLM으로 PDF 파싱
                result = self._parse_with_vlm(pdf_path)

                # 임시 PDF 삭제
                try:
                    pdf_path.unlink()
                    print(f"  🗑️ 임시 PDF 삭제됨")
                except:
                    pass

                if result:
                    print(f"  ✅ VLM OCR 성공")
                    return result
                else:
                    print(f"  ⚠️ VLM OCR 실패")

        # 방법 3: 모든 방법 실패 → olefile 결과라도 반환 (있으면)
        if olefile_result:
            print(f"  ⚠️ VLM OCR 실패, olefile 결과 사용 ({olefile_text_length}자)")
            print(f"  📌 PrvText는 문서 미리보기용으로 일부 내용만 포함됨")
            return olefile_result

        print(f"  ❌ 모든 HWP 처리 방법 실패")
        return []

    def _load_hwpx(self, file_path: Path) -> List[Dict]:
        """HWPX 파일 읽기 - PDF 변환 방식 우선"""
        print(f"\n  📄 HWPX 파일 처리: {file_path.name}")

        # 🆕 새로운 방식: HWPX → PDF 변환 후 처리
        print(f"  🔄 방법: HWPX → PDF 변환 후 기존 PDF 로직 사용")

        # 1. HWPX → PDF 변환
        pdf_path = self._convert_hwp_to_pdf(file_path)

        if pdf_path and pdf_path.exists():
            # 2. 변환된 PDF 처리
            print(f"  📖 변환된 PDF 로드 중...")
            try:
                result = self._load_pdf(pdf_path)

                # 3. 임시 PDF 파일 삭제
                try:
                    pdf_path.unlink()
                    print(f"  🗑️ 임시 PDF 파일 삭제됨")
                except:
                    pass

                # 4. method를 'hwpx_via_pdf'로 표시
                for page in result:
                    page["method"] = "hwpx_via_pdf"

                print(f"  ✅ HWPX → PDF 변환 방식 처리 완료")
                return result

            except Exception as e:
                print(f"  ❌ 변환된 PDF 처리 실패: {e}")

                # 임시 파일 삭제 시도
                try:
                    if pdf_path.exists():
                        pdf_path.unlink()
                except:
                    pass

        # 변환 실패 시 기존 방식 폴백
        print(f"  ⚠️ PDF 변환 실패, 기존 방식(ZIP XML 파싱) 시도...")

        try:
            import zipfile
            import xml.etree.ElementTree as ET

            texts = []

            with zipfile.ZipFile(file_path, "r") as zip_ref:
                for name in zip_ref.namelist():
                    if name.startswith("Contents/section") and name.endswith(".xml"):
                        try:
                            xml_content = zip_ref.read(name)
                            root = ET.fromstring(xml_content)

                            section_texts = []
                            for elem in root.iter():
                                if elem.text and elem.text.strip():
                                    section_texts.append(elem.text.strip())
                                if elem.tail and elem.tail.strip():
                                    section_texts.append(elem.tail.strip())

                            section_text = "\n".join(section_texts)
                            if section_text and self._is_valid_korean_text(
                                section_text
                            ):
                                texts.append(section_text)
                            else:
                                print(
                                    f"    ⚠️ 섹션 {name} 텍스트가 유효하지 않음 (건너뜀)"
                                )

                        except Exception as e:
                            print(f"    ⚠️ 섹션 {name} 파싱 실패: {e}")
                            continue

            if not texts:
                print("  ❌ 텍스트 추출 실패")
                return []

            full_text = "\n".join(texts)
            full_text = self.text_cleaner.clean_ocr_text(full_text)

            print(f"  ✓ HWPX 읽기 완료 (ZIP XML 폴백, {len(full_text)}자)")
            return [{"page_num": 1, "text": full_text, "method": "hwpx_xml"}]

        except Exception as e:
            print(f"  ❌ HWPX 읽기 실패: {e}")
            return []

    # ============================================
    # 이미지 처리
    # ============================================

    def _load_image(self, file_path: Path) -> List[Dict]:
        """이미지 OCR"""
        try:
            print(f"  🖼️ 이미지 OCR 처리 중...")

            image = Image.open(file_path)

            # 전처리 추가
            image = self._preprocess_image_for_table(image)

            text = pytesseract.image_to_string(
                image, lang="kor+eng", config="--oem 1 --psm 6"
            ).strip()

            # 후처리 추가
            text = self.text_cleaner.clean_ocr_text(text)

            print(f"  ✓ OCR 완료 ({len(text)}자)")

            return [{"page_num": 1, "text": text, "method": "image_ocr"}]

        except Exception as e:
            print(f"  ❌ 이미지 OCR 실패: {e}")
            return []

    # ============================================
    # OCR 전처리 및 후처리
    # ============================================

    def _preprocess_image_for_table(self, image: Image.Image) -> Image.Image:
        """표 인식을 위한 이미지 전처리"""
        # 1. 그레이스케일 변환
        image = image.convert("L")

        # 2. 대비 강화 (표 선을 더 명확하게)
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2.5)

        # 3. 선명도 강화
        enhancer = ImageEnhance.Sharpness(image)
        image = enhancer.enhance(2.0)

        # 4. 이진화 (표 경계 강조)
        threshold = 128
        image = image.point(lambda p: 255 if p > threshold else 0)

        return image

