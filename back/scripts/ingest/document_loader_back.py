"""
범용 문서 로더 (Universal Document Loader)
모든 파일 형식을 자동 감지하고 텍스트 추출
"""

from pathlib import Path
from typing import List, Dict
import platform

# PDF
import PyPDF2
from pdf2image import convert_from_path
import pytesseract
from PIL import Image

# Office 문서
try:
    import docx

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import pandas as pd

    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

# Tesseract/Poppler 경로 (Windows)
if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = (
        r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    )

    import glob

    poppler_search = glob.glob(r"C:\Program Files\poppler-*\Library\bin")
    POPPLER_PATH = (
        poppler_search[0]
        if poppler_search
        else r"C:\Program Files\poppler-25.07.0\Library\bin"
    )
else:
    POPPLER_PATH = None


class UniversalDocumentLoader:
    """범용 문서 로더 - 모든 파일 형식 자동 처리"""

    def __init__(self, config):
        self.config = config
        self.ocr_dpi = getattr(config, "ocr_dpi", 300)
        self._print_capabilities()

    def _print_capabilities(self):
        """지원 형식 출력"""
        print("\n📋 문서 로더 지원 형식:")
        print(f"  - PDF: ✓ (텍스트 + OCR)")
        print(f"  - TXT: ✓")
        print(
            f"  - Word (.docx): {'✓' if DOCX_AVAILABLE else '✗ (pip install python-docx)'}"
        )
        print(
            f"  - PowerPoint (.pptx): {'✓' if PPTX_AVAILABLE else '✗ (pip install python-pptx)'}"
        )
        print(
            f"  - Excel (.xlsx): {'✓' if OPENPYXL_AVAILABLE else '✗ (pip install openpyxl)'}"
        )
        print(f"  - CSV: {'✓' if PANDAS_AVAILABLE else '✗ (pip install pandas)'}")
        print(f"  - 이미지 (.jpg/.png): ✓ (OCR)\n")

    def load(self, file_path: Path) -> List[Dict]:
        """
        파일 형식 자동 감지 및 텍스트 추출

        Returns:
            List[Dict]: [{'page_num': int, 'text': str, 'method': str}, ...]
        """
        suffix = file_path.suffix.lower()

        # 파일 형식별 라우팅
        if suffix == ".pdf":
            return self._load_pdf(file_path)
        elif suffix == ".txt":
            return self._load_txt(file_path)
        elif suffix in [".docx", ".doc"]:
            return self._load_docx(file_path)
        elif suffix in [".pptx", ".ppt"]:
            return self._load_pptx(file_path)
        elif suffix in [".xlsx", ".xls"]:
            return self._load_excel(file_path)
        elif suffix == ".csv":
            return self._load_csv(file_path)
        elif suffix in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
            return self._load_image(file_path)
        else:
            print(f"  ⚠️ 지원하지 않는 형식: {suffix}")
            return []

    # ============================================
    # PDF 처리
    # ============================================

    def _load_pdf(self, file_path: Path) -> List[Dict]:
        """PDF 텍스트 추출 (텍스트 우선, 실패시 OCR)"""
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)

                print(f"  📑 PDF {total_pages}페이지")

                # 첫 페이지로 스캔 여부 판단
                first_text = reader.pages[0].extract_text().strip()

                if len(first_text) < 50:
                    print("  ⚠️ 스캔본 감지 → 전체 OCR")
                    return self._ocr_pdf(file_path)

                # 텍스트 기반 PDF
                print("  ✓ 텍스트 기반 PDF")
                pages_data = []

                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text().strip()

                    if len(text) < 50:
                        print(f"    페이지 {page_num}: OCR 적용")
                        text = self._ocr_pdf_page(file_path, page_num)
                        method = "pdf_ocr"
                    else:
                        method = "pdf_text"

                    pages_data.append(
                        {"page_num": page_num, "text": text, "method": method}
                    )

                return pages_data

        except Exception as e:
            print(f"  ❌ PDF 읽기 실패: {e}")
            print("  → OCR 모드로 전환")
            return self._ocr_pdf(file_path)

    def _ocr_pdf(self, file_path: Path) -> List[Dict]:
        """PDF 전체 OCR"""
        try:
            if platform.system() == "Windows":
                images = convert_from_path(
                    file_path, dpi=self.ocr_dpi, poppler_path=POPPLER_PATH
                )
            else:
                images = convert_from_path(file_path, dpi=self.ocr_dpi)

            print(f"  🔍 {len(images)}페이지 OCR 처리 중...")

            pages_data = []
            for page_num, image in enumerate(images, 1):
                print(f"    페이지 {page_num}/{len(images)}...", end=" ")

                try:
                    text = pytesseract.image_to_string(
                        image, lang="kor+eng", config="--oem 3 --psm 6"
                    ).strip()

                    print(f"✓ ({len(text)}자)")

                    pages_data.append(
                        {"page_num": page_num, "text": text, "method": "pdf_ocr"}
                    )
                except Exception as e:
                    print(f"✗ 실패: {e}")
                    pages_data.append(
                        {"page_num": page_num, "text": "", "method": "ocr_failed"}
                    )

            return pages_data

        except Exception as e:
            print(f"  ❌ OCR 실패: {e}")
            return []

    def _ocr_pdf_page(self, file_path: Path, page_num: int) -> str:
        """PDF 특정 페이지만 OCR"""
        try:
            if platform.system() == "Windows":
                images = convert_from_path(
                    file_path,
                    first_page=page_num,
                    last_page=page_num,
                    dpi=self.ocr_dpi,
                    poppler_path=POPPLER_PATH,
                )
            else:
                images = convert_from_path(
                    file_path, first_page=page_num, last_page=page_num, dpi=self.ocr_dpi
                )

            if images:
                text = pytesseract.image_to_string(
                    images[0], lang="kor+eng", config="--oem 3 --psm 6"
                ).strip()
                return text
            return ""

        except Exception as e:
            print(f" (OCR 실패: {e})")
            return ""

    # ============================================
    # TXT 처리
    # ============================================

    def _load_txt(self, file_path: Path) -> List[Dict]:
        """TXT 파일 읽기"""
        try:
            # 인코딩 자동 감지
            encodings = ["utf-8", "cp949", "euc-kr", "latin-1"]
            text = None

            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        text = f.read()
                    print(f"  ✓ TXT 읽기 성공 ({encoding})")
                    break
                except UnicodeDecodeError:
                    continue

            if text is None:
                print("  ❌ 인코딩 실패")
                return []

            return [{"page_num": 1, "text": text, "method": "txt"}]

        except Exception as e:
            print(f"  ❌ TXT 읽기 실패: {e}")
            return []

    # ============================================
    # Word 처리
    # ============================================

    def _load_docx(self, file_path: Path) -> List[Dict]:
        """Word 문서 읽기"""
        if not DOCX_AVAILABLE:
            print("  ❌ python-docx 미설치")
            return []

        try:
            doc = docx.Document(file_path)

            # 단락별 텍스트 추출
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

            # 표 추출
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        tables_text.append(row_text)

            # 병합
            all_text = "\n".join(paragraphs)
            if tables_text:
                all_text += "\n\n[표 데이터]\n" + "\n".join(tables_text)

            print(f"  ✓ Word 읽기 완료 ({len(all_text)}자)")

            return [{"page_num": 1, "text": all_text, "method": "docx"}]

        except Exception as e:
            print(f"  ❌ Word 읽기 실패: {e}")
            return []

    # ============================================
    # PowerPoint 처리
    # ============================================

    def _load_pptx(self, file_path: Path) -> List[Dict]:
        """PowerPoint 읽기"""
        if not PPTX_AVAILABLE:
            print("  ❌ python-pptx 미설치")
            return []

        try:
            prs = Presentation(file_path)
            pages_data = []

            print(f"  📊 PowerPoint {len(prs.slides)}슬라이드")

            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())

                slide_text = "\n".join(texts)

                pages_data.append(
                    {"page_num": slide_num, "text": slide_text, "method": "pptx"}
                )

            print(f"  ✓ PowerPoint 읽기 완료")
            return pages_data

        except Exception as e:
            print(f"  ❌ PowerPoint 읽기 실패: {e}")
            return []

    # ============================================
    # Excel 처리
    # ============================================

    def _load_excel(self, file_path: Path) -> List[Dict]:
        """Excel 읽기"""
        if not OPENPYXL_AVAILABLE:
            print("  ❌ openpyxl 미설치")
            return []

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            pages_data = []

            print(f"  📊 Excel {len(wb.sheetnames)}시트")

            for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
                sheet = wb[sheet_name]

                rows_text = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip():
                        rows_text.append(row_text)

                sheet_text = f"[시트: {sheet_name}]\n" + "\n".join(rows_text)

                pages_data.append(
                    {"page_num": sheet_num, "text": sheet_text, "method": "excel"}
                )

            print(f"  ✓ Excel 읽기 완료")
            return pages_data

        except Exception as e:
            print(f"  ❌ Excel 읽기 실패: {e}")
            return []

    # ============================================
    # CSV 처리
    # ============================================

    def _load_csv(self, file_path: Path) -> List[Dict]:
        """CSV 읽기"""
        if not PANDAS_AVAILABLE:
            print("  ❌ pandas 미설치")
            return []

        try:
            # 인코딩 자동 감지
            encodings = ["utf-8", "cp949", "euc-kr"]
            df = None

            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    print(f"  ✓ CSV 읽기 성공 ({encoding})")
                    break
                except:
                    continue

            if df is None:
                print("  ❌ CSV 인코딩 실패")
                return []

            # 데이터프레임을 텍스트로 변환
            csv_text = df.to_string(index=False)

            print(f"  ✓ CSV 읽기 완료 ({len(df)}행)")

            return [{"page_num": 1, "text": csv_text, "method": "csv"}]

        except Exception as e:
            print(f"  ❌ CSV 읽기 실패: {e}")
            return []

    # ============================================
    # 이미지 처리
    # ============================================

    def _load_image(self, file_path: Path) -> List[Dict]:
        """이미지 OCR"""
        try:
            print(f"  🖼️ 이미지 OCR 처리 중...")

            image = Image.open(file_path)

            text = pytesseract.image_to_string(
                image, lang="kor+eng", config="--oem 3 --psm 6"
            ).strip()

            print(f"  ✓ OCR 완료 ({len(text)}자)")

            return [{"page_num": 1, "text": text, "method": "image_ocr"}]

        except Exception as e:
            print(f"  ❌ 이미지 OCR 실패: {e}")
            return []
